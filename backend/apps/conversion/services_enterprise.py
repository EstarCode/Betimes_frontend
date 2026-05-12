"""
Betimes - Enterprise Document Conversion Services
Complete document format conversion toolkit
Windows-compatible without external dependencies
"""

import os
import logging
import subprocess
from typing import Optional, Dict
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
import pdfplumber
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.utils import ImageReader
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors

from config.exceptions import ConversionError

logger = logging.getLogger(__name__)


class EnterpriseConversionService:
    """Enterprise-grade document conversion service"""
    
    def __init__(self):
        self.temp_dir = '/tmp/docforge'
        os.makedirs(self.temp_dir, exist_ok=True)
    
    # ==================== WORD CONVERSIONS ====================
    
    def word_to_pdf(self, input_path: str, output_path: str) -> Dict:
        """
        Convert Word (.docx) to PDF
        Uses reportlab for conversion (LibreOffice optional)
        """
        try:
            logger.info(f"Converting Word to PDF: {input_path}")
            
            # Try using LibreOffice if available (best quality)
            try:
                subprocess.run([
                    'soffice',
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', os.path.dirname(output_path),
                    input_path
                ], check=True, timeout=60)
                
                # Rename output if needed
                expected_output = os.path.join(
                    os.path.dirname(output_path),
                    Path(input_path).stem + '.pdf'
                )
                if expected_output != output_path and os.path.exists(expected_output):
                    os.rename(expected_output, output_path)
                    
            except (subprocess.CalledProcessError, FileNotFoundError):
                # Fallback to reportlab
                logger.warning("LibreOffice not available, using reportlab")
                
                doc = Document(input_path)
                c = canvas.Canvas(output_path, pagesize=letter)
                width, height = letter
                
                y_position = height - 50
                for para in doc.paragraphs:
                    if y_position < 50:
                        c.showPage()
                        y_position = height - 50
                    
                    text = para.text[:100]  # Limit line length
                    c.drawString(50, y_position, text)
                    y_position -= 20
                
                c.save()
            
            logger.info(f"Word to PDF conversion completed: {output_path}")
            
            return {
                'success': True,
                'output_size': os.path.getsize(output_path),
                'format': 'PDF'
            }
            
        except Exception as e:
            logger.exception(f"Word to PDF conversion failed: {str(e)}")
            raise ConversionError(f"Word to PDF conversion failed: {str(e)}")
    
    def pdf_to_word(self, input_path: str, output_path: str) -> Dict:
        """
        Convert PDF to Word (.docx)
        Extracts text and formatting
        """
        try:
            logger.info(f"Converting PDF to Word: {input_path}")
            
            # Create new Word document
            doc = Document()
            
            # Extract text from PDF
            with pdfplumber.open(input_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    # Add page heading
                    doc.add_heading(f'Page {page_num}', level=2)
                    
                    # Extract text
                    text = page.extract_text()
                    if text:
                        doc.add_paragraph(text)
                    
                    # Extract tables
                    tables = page.extract_tables()
                    for table_data in tables:
                        if table_data:
                            # Add table to Word
                            word_table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
                            for i, row in enumerate(table_data):
                                for j, cell in enumerate(row):
                                    word_table.rows[i].cells[j].text = str(cell) if cell else ''
                    
                    # Page break
                    if page_num < len(pdf.pages):
                        doc.add_page_break()
            
            doc.save(output_path)
            
            logger.info(f"PDF to Word conversion completed: {output_path}")
            
            return {
                'success': True,
                'output_size': os.path.getsize(output_path),
                'format': 'DOCX'
            }
            
        except Exception as e:
            logger.exception(f"PDF to Word conversion failed: {str(e)}")
            raise ConversionError(f"PDF to Word conversion failed: {str(e)}")
    
    def word_to_excel(self, input_path: str, output_path: str) -> Dict:
        """
        Convert Word tables to Excel
        Extracts tables from Word document
        """
        try:
            logger.info(f"Converting Word to Excel: {input_path}")
            
            doc = Document(input_path)
            wb = Workbook()
            ws = wb.active
            ws.title = "Extracted Data"
            
            row_num = 1
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    for col_num, cell in enumerate(row.cells, 1):
                        ws.cell(row=row_num, column=col_num, value=cell.text)
                    row_num += 1
                row_num += 1  # Add space between tables
            
            # If no tables, extract paragraphs
            if len(doc.tables) == 0:
                for para in doc.paragraphs:
                    ws.cell(row=row_num, column=1, value=para.text)
                    row_num += 1
            
            wb.save(output_path)
            
            logger.info(f"Word to Excel conversion completed: {output_path}")
            
            return {
                'success': True,
                'tables_extracted': len(doc.tables),
                'output_size': os.path.getsize(output_path)
            }
            
        except Exception as e:
            logger.exception(f"Word to Excel conversion failed: {str(e)}")
            raise ConversionError(f"Word to Excel conversion failed: {str(e)}")
    
    # ==================== EXCEL CONVERSIONS ====================
    
    def excel_to_pdf(self, input_path: str, output_path: str) -> Dict:
        """Convert Excel to PDF"""
        try:
            logger.info(f"Converting Excel to PDF: {input_path}")
            
            # Try LibreOffice first
            try:
                subprocess.run([
                    'soffice',
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', os.path.dirname(output_path),
                    input_path
                ], check=True, timeout=60)
                
                expected_output = os.path.join(
                    os.path.dirname(output_path),
                    Path(input_path).stem + '.pdf'
                )
                if expected_output != output_path and os.path.exists(expected_output):
                    os.rename(expected_output, output_path)
                    
            except (subprocess.CalledProcessError, FileNotFoundError):
                # Fallback: Create PDF from Excel data
                logger.warning("LibreOffice not available, using reportlab")
                
                wb = load_workbook(input_path)
                ws = wb.active
                
                # Create PDF
                doc = SimpleDocTemplate(output_path, pagesize=A4)
                elements = []
                
                # Convert Excel data to table
                data = []
                for row in ws.iter_rows(values_only=True):
                    data.append([str(cell) if cell is not None else '' for cell in row])
                
                if data:
                    table = Table(data)
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), 10),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black)
                    ]))
                    elements.append(table)
                
                doc.build(elements)
            
            logger.info(f"Excel to PDF conversion completed: {output_path}")
            
            return {
                'success': True,
                'output_size': os.path.getsize(output_path)
            }
            
        except Exception as e:
            logger.exception(f"Excel to PDF conversion failed: {str(e)}")
            raise ConversionError(f"Excel to PDF conversion failed: {str(e)}")
    
    def pdf_to_excel(self, input_path: str, output_path: str) -> Dict:
        """
        Convert PDF to Excel
        Extracts tables from PDF
        """
        try:
            logger.info(f"Converting PDF to Excel: {input_path}")
            
            wb = Workbook()
            
            with pdfplumber.open(input_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    # Create sheet for each page
                    if page_num == 1:
                        ws = wb.active
                        ws.title = f"Page {page_num}"
                    else:
                        ws = wb.create_sheet(title=f"Page {page_num}")
                    
                    # Extract tables
                    tables = page.extract_tables()
                    
                    row_num = 1
                    for table in tables:
                        for row_data in table:
                            for col_num, cell in enumerate(row_data, 1):
                                ws.cell(row=row_num, column=col_num, value=cell)
                            row_num += 1
                        row_num += 1  # Space between tables
                    
                    # If no tables, extract text
                    if not tables:
                        text = page.extract_text()
                        if text:
                            for line in text.split('\n'):
                                ws.cell(row=row_num, column=1, value=line)
                                row_num += 1
            
            wb.save(output_path)
            
            logger.info(f"PDF to Excel conversion completed: {output_path}")
            
            return {
                'success': True,
                'sheets_created': len(wb.sheetnames),
                'output_size': os.path.getsize(output_path)
            }
            
        except Exception as e:
            logger.exception(f"PDF to Excel conversion failed: {str(e)}")
            raise ConversionError(f"PDF to Excel conversion failed: {str(e)}")
    
    # ==================== POWERPOINT CONVERSIONS ====================
    
    def ppt_to_pdf(self, input_path: str, output_path: str) -> Dict:
        """Convert PowerPoint to PDF"""
        try:
            logger.info(f"Converting PowerPoint to PDF: {input_path}")
            
            # Use LibreOffice
            subprocess.run([
                'soffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', os.path.dirname(output_path),
                input_path
            ], check=True, timeout=120)
            
            expected_output = os.path.join(
                os.path.dirname(output_path),
                Path(input_path).stem + '.pdf'
            )
            if expected_output != output_path and os.path.exists(expected_output):
                os.rename(expected_output, output_path)
            
            logger.info(f"PowerPoint to PDF conversion completed: {output_path}")
            
            return {
                'success': True,
                'output_size': os.path.getsize(output_path)
            }
            
        except Exception as e:
            logger.exception(f"PowerPoint to PDF conversion failed: {str(e)}")
            raise ConversionError(f"PowerPoint to PDF conversion failed: {str(e)}")
    
    def pdf_to_ppt(self, input_path: str, output_path: str) -> Dict:
        """
        Convert PDF to PowerPoint
        Creates one slide per page (requires pdf2image and poppler)
        """
        try:
            logger.info(f"Converting PDF to PowerPoint: {input_path}")
            
            # Note: This requires pdf2image library and poppler
            try:
                from pdf2image import convert_from_path
                
                prs = Presentation()
                
                # Convert PDF pages to images
                images = convert_from_path(input_path, dpi=150)
                
                for img in images:
                    # Add blank slide
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                    
                    # Save image temporarily
                    img_path = os.path.join(self.temp_dir, 'temp_slide.png')
                    img.save(img_path, 'PNG')
                    
                    # Add image to slide
                    slide.shapes.add_picture(
                        img_path,
                        left=0,
                        top=0,
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                    
                    # Clean up temp image
                    if os.path.exists(img_path):
                        os.remove(img_path)
                
                prs.save(output_path)
                
                logger.info(f"PDF to PowerPoint conversion completed: {output_path}")
                
                return {
                    'success': True,
                    'slides_created': len(images),
                    'output_size': os.path.getsize(output_path)
                }
                
            except ImportError:
                logger.warning("pdf2image not available")
                raise ConversionError("PDF to PowerPoint conversion requires pdf2image library and poppler")
            
        except Exception as e:
            logger.exception(f"PDF to PowerPoint conversion failed: {str(e)}")
            raise ConversionError(f"PDF to PowerPoint conversion failed: {str(e)}")
    
    # ==================== IMAGE CONVERSIONS ====================
    
    def image_to_pdf(self, input_path: str, output_path: str) -> Dict:
        """Convert image to PDF"""
        try:
            logger.info(f"Converting image to PDF: {input_path}")
            
            img = Image.open(input_path)
            
            # Convert to RGB if necessary
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            # Save as PDF
            img.save(output_path, 'PDF', resolution=100.0)
            
            logger.info(f"Image to PDF conversion completed: {output_path}")
            
            return {
                'success': True,
                'image_size': img.size,
                'output_size': os.path.getsize(output_path)
            }
            
        except Exception as e:
            logger.exception(f"Image to PDF conversion failed: {str(e)}")
            raise ConversionError(f"Image to PDF conversion failed: {str(e)}")
    
    def pdf_to_images(self, input_path: str, output_dir: str, 
                     format: str = 'PNG', dpi: int = 200) -> Dict:
        """Convert PDF pages to images (requires pdf2image and poppler)"""
        try:
            logger.info(f"Converting PDF to images: {input_path}")
            
            os.makedirs(output_dir, exist_ok=True)
            
            # Note: This requires pdf2image library and poppler
            try:
                from pdf2image import convert_from_path
                
                # Convert PDF to images
                images = convert_from_path(input_path, dpi=dpi)
                
                output_files = []
                for i, img in enumerate(images, 1):
                    output_file = os.path.join(output_dir, f"page_{i}.{format.lower()}")
                    img.save(output_file, format)
                    output_files.append(output_file)
                
                logger.info(f"PDF to images conversion completed: {len(images)} pages")
                
                return {
                    'success': True,
                    'pages_converted': len(images),
                    'output_files': output_files
                }
                
            except ImportError:
                logger.warning("pdf2image not available")
                raise ConversionError("PDF to images conversion requires pdf2image library and poppler")
            
        except Exception as e:
            logger.exception(f"PDF to images conversion failed: {str(e)}")
            raise ConversionError(f"PDF to images conversion failed: {str(e)}")
    
    # ==================== OCR ====================
    
    def ocr_pdf(self, input_path: str, output_path: str, language: str = 'eng') -> Dict:
        """
        Perform OCR on scanned PDF
        Requires pytesseract and pdf2image libraries
        """
        try:
            logger.info(f"Performing OCR on PDF: {input_path}")
            
            # Note: This requires pytesseract, pdf2image, and poppler
            try:
                from pdf2image import convert_from_path
                import pytesseract
                
                # Convert PDF to images
                images = convert_from_path(input_path, dpi=300)
                
                # Create new PDF with OCR text
                doc = fitz.open()
                
                for img in images:
                    # Perform OCR
                    text = pytesseract.image_to_string(img, lang=language)
                    
                    # Create new page
                    page = doc.new_page(width=img.width, height=img.height)
                    
                    # Add image
                    import io
                    img_bytes = io.BytesIO()
                    img.save(img_bytes, format='PNG')
                    img_bytes.seek(0)
                    
                    page.insert_image(page.rect, stream=img_bytes.read())
                    
                    # Add invisible text layer
                    page.insert_text((10, 10), text, fontsize=1, color=(1, 1, 1))
                
                doc.save(output_path)
                doc.close()
                
                logger.info(f"OCR completed: {output_path}")
                
                return {
                    'success': True,
                    'pages_processed': len(images),
                    'output_size': os.path.getsize(output_path)
                }
                
            except ImportError as ie:
                logger.warning(f"OCR dependencies not available: {str(ie)}")
                raise ConversionError("OCR requires pytesseract, pdf2image libraries and Tesseract OCR engine")
            
        except Exception as e:
            logger.exception(f"OCR failed: {str(e)}")
            raise ConversionError(f"OCR failed: {str(e)}")
